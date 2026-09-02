"""Pre-submission preflight. Fails loudly rather than reassuring quietly."""
import json, os, re, sys
import pandas as pd

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
D = os.path.join(BASE, "SUBMISSION")
ok, warn, fail = [], [], []


def chk(cond, good, bad, hard=True):
    (ok if cond else (fail if hard else warn)).append(good if cond else bad)


# ---- files present -------------------------------------------------------
need = ["WRITEUP.md", "METHODOLOGY.md", "README.md",
        "aisehack_round3_sar_yield_forecast_EXECUTED.ipynb",
        "aisehack_round3_sar_yield_forecast.ipynb",
        "Team8bit_Round3_GoaFinals.pptx", "Team8bit_Round3_results.xlsx",
        "plot_level_yield_forecast.csv", "village_level_yield_forecast.csv",
        "crop_mix_scenarios.csv", "validation_report.json"]
for f in need:
    p = os.path.join(D, f)
    chk(os.path.exists(p) and os.path.getsize(p) > 200,
        f"{f} present ({os.path.getsize(p)/1024:.1f} KB)" if os.path.exists(p) else "",
        f"MISSING OR EMPTY: {f}")

figs = sorted(os.listdir(os.path.join(D, "figures")))
chk(len(figs) >= 9, f"{len(figs)} figures", f"only {len(figs)} figures")
chk("01_cover.png" in figs, "cover image present", "NO COVER IMAGE")
chk("00_architecture.png" in figs, "architecture diagram present", "no architecture diagram")

# ---- writeup -------------------------------------------------------------
t = open(os.path.join(D, "WRITEUP.md"), encoding="utf-8").read()
w = len(re.sub(r"\|", "", t).split())
chk(w <= 2000, f"writeup {w} words (limit 2000)", f"WRITEUP OVER LIMIT: {w} words")
secs = [l for l in t.splitlines() if l.startswith("## ")]
nums = [int(re.match(r"## (\d+)\.", s).group(1)) for s in secs if re.match(r"## (\d+)\.", s)]
chk(nums == list(range(1, len(nums) + 1)),
    f"writeup sections numbered 1..{len(nums)}", f"SECTION NUMBERING BROKEN: {nums}")
chk("Team 8bit" in t and "Harsh Thummar" in t and "Viraj Suhagiya" in t,
    "team and authors named", "AUTHORS MISSING")

# ---- plot / village tables ----------------------------------------------
p = pd.read_csv(os.path.join(D, "plot_level_yield_forecast.csv"))
v = pd.read_csv(os.path.join(D, "village_level_yield_forecast.csv"))
chk(len(p) == 966, "966 plot rows", f"WRONG ROW COUNT: {len(p)}")
chk(p.farm_id.nunique() == 966, "farm_id unique", "DUPLICATE farm_id")
chk(p.final_yield_forecast_kg_ha.notna().all(), "no missing forecasts", "MISSING FORECASTS")
chk((p.final_yield_forecast_kg_ha > 0).all(), "all forecasts positive", "NON-POSITIVE FORECAST")
bad = ((p.yield_p10_kg_ha > p.final_yield_forecast_kg_ha) |
       (p.yield_p90_kg_ha < p.final_yield_forecast_kg_ha)).sum()
chk(bad == 0, "every plot interval brackets its central value",
    f"{bad} PLOT INTERVALS DO NOT BRACKET")
chk(set(p.crop_type.unique()) <= {"Rice", "Cotton", "Maize", "Bajra", "Groundnut"},
    "only the five permitted crops", "UNEXPECTED CROP LABEL")
chk(abs(p.area_ha.sum() - 447.54) < 0.1, f"area {p.area_ha.sum():.2f} ha", "AREA MISMATCH")

for _, r in v.iterrows():
    chk(r.yield_p10_kg_ha <= r.yield_forecast_kg_ha <= r.yield_p90_kg_ha,
        f"{r.crop_type} village interval brackets", f"{r.crop_type} INTERVAL BROKEN")
tot_p = (p.final_yield_forecast_kg_ha * p.area_ha).sum() / 1000
chk(abs(tot_p - v.production_t.sum()) < 0.5,
    f"plot and village totals agree ({tot_p:.1f} t)", "PLOT/VILLAGE TOTALS DISAGREE")

# ---- executed notebook ---------------------------------------------------
nb = json.load(open(os.path.join(D, "aisehack_round3_sar_yield_forecast_EXECUTED.ipynb"),
                    encoding="utf-8"))
code = [c for c in nb["cells"] if c["cell_type"] == "code"]
withio = [c for c in code if c.get("outputs")]
errs = [c for c in code for o in c.get("outputs", []) if o.get("output_type") == "error"]
chk(len(errs) == 0, "no error outputs in the executed notebook",
    f"{len(errs)} CELLS ERRORED IN THE NOTEBOOK")
chk(len(withio) >= len(code) - 2,
    f"{len(withio)}/{len(code)} code cells carry outputs", "notebook outputs missing", hard=False)
alltxt = json.dumps(nb)
chk("all assertions passed" in alltxt, "notebook self-check passed",
    "NOTEBOOK ASSERTIONS DID NOT RUN")

# ---- pptx ----------------------------------------------------------------
try:
    from pptx import Presentation
    pr = Presentation(os.path.join(D, "Team8bit_Round3_GoaFinals.pptx"))
    n = len(pr.slides)
    notes = sum(1 for s in pr.slides
                if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    chk(n >= 14, f"deck has {n} slides", f"deck only has {n} slides")
    chk(notes == n, f"speaker notes on all {n} slides", f"notes on only {notes}/{n} slides",
        hard=False)
    lit = sum(1 for s in pr.slides for sh in s.shapes
              if sh.has_text_frame and "**" in sh.text_frame.text)
    chk(lit == 0, "no literal markdown markers on slides", f"{lit} SLIDES SHOW ** MARKERS")
except Exception as e:
    fail.append(f"PPTX WILL NOT OPEN: {e}")

# ---- validation report ---------------------------------------------------
rep = json.load(open(os.path.join(D, "validation_report.json"), encoding="utf-8"))
chk(rep["site"]["n_plots"] == 966, "report says 966 plots", "report plot count wrong")
chk("limitations" in rep and len(rep["limitations"]) >= 5,
    f"{len(rep.get('limitations', []))} limitations documented", "limitations missing")

sc = pd.read_csv(os.path.join(D, "crop_mix_scenarios.csv"))
chk(len(sc) == 10, "both crop-mix scenarios present (5 crops x 2)",
    f"scenario file has {len(sc)} rows")

# ---- report --------------------------------------------------------------
print("=" * 72)
print(f"PASSED {len(ok)}")
for x in ok:
    print("  ok   ", x)
if warn:
    print(f"\nWARNINGS {len(warn)}")
    for x in warn:
        print("  warn ", x)
print("\n" + "=" * 72)
if fail:
    print(f"BLOCKERS {len(fail)}")
    for x in fail:
        print("  FAIL ", x)
    sys.exit(1)
print("NO BLOCKERS — deliverables are internally consistent and complete.")
