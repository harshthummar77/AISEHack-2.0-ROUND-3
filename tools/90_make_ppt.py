"""Goa Finals deck — 10 minutes, presentation-friendly.

Design rules applied throughout:
  * one idea per slide, headline carries the message
  * at most three supporting lines; detail lives in the speaker notes
  * every number on a slide is one we can defend if challenged
"""
import os, re
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(BASE, "SUBMISSION", "figures")
OUT = os.path.join(BASE, "pipeline", "out")
DEL = os.path.join(BASE, "SUBMISSION")

INK = RGBColor(0x14, 0x1B, 0x24)
MUTE = RGBColor(0x5B, 0x68, 0x76)
ACC = RGBColor(0x15, 0x67, 0x8A)
GOOD = RGBColor(0x2C, 0x72, 0x54)
WARN = RGBColor(0x8F, 0x5F, 0x14)
BAD = RGBColor(0x9C, 0x35, 0x29)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF1, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

vil = pd.read_csv(os.path.join(OUT, "village_yield_forecast.csv"))
_TOK = re.compile(r"(\*\*.+?\*\*|\*.+?\*)")


def slide(notes=None):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def tb(s, x, y, w, h, text, size=18, bold=False, color=INK,
       align=PP_ALIGN.LEFT, space_after=8, line=1.2):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line
        for part in _TOK.split(ln):
            if not part:
                continue
            b, it = bold, False
            if part.startswith("**") and part.endswith("**") and len(part) > 4:
                part, b = part[2:-2], True
            elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                part, it = part[1:-1], True
            r = p.add_run(); r.text = part
            r.font.size = Pt(size); r.font.bold = b; r.font.italic = it
            r.font.color.rgb = color; r.font.name = "Segoe UI"
    return box


def title(s, t, sub=None):
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = BAND
    bar.line.fill.background(); bar.shadow.inherit = False
    tb(s, 0.6, 0.2, 12.2, 0.6, t, size=30, bold=True)
    if sub:
        tb(s, 0.6, 0.76, 12.2, 0.34, sub, size=14, color=MUTE)


def pic(s, name, x, y, w=None, h=None):
    p = os.path.join(FIG, name)
    if not os.path.exists(p):
        return None
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(p, Inches(x), Inches(y), **kw)


def stat(s, x, y, w, big, label, color=ACC, bs=40):
    tb(s, x, y, w, 0.75, big, size=bs, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(s, x, y + 0.78, w, 0.8, label, size=12, color=MUTE, align=PP_ALIGN.CENTER)


def table(s, df, x, y, w, h, fs=14, hdr=13, col_w=None, hi=None):
    r, c = df.shape[0] + 1, df.shape[1]
    t = s.shapes.add_table(r, c, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_w:
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            t.columns[i].width = Emu(int(Inches(w) * cw / tot))
    def put(cell, text, size, bold, colour, align):
        # an empty string creates no run, so style the run only when there is one
        cell.text = str(text)
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        if not p.runs:
            p.add_run().text = ""
        r = p.runs[0]
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = colour; r.font.name = "Segoe UI"

    for j, cn in enumerate(df.columns):
        cell = t.cell(0, j)
        put(cell, cn, hdr, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
        cell.fill.solid(); cell.fill.fore_color.rgb = ACC
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i in range(df.shape[0]):
        for j in range(c):
            cell = t.cell(i + 1, j)
            put(cell, df.iloc[i, j], fs, hi is not None and i in hi, INK,
                PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if i % 2 == 0 else BAND
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return t


# ============================================================ 1 title
s = slide("Good morning. Team 8bit. Sokhda village, Vadodara — 966 plots, 447 hectares, "
          "six Capella X-band passes across kharif 2025. Our job this round is a final "
          "at-harvest forecast, not a snapshot. Ten minutes, and I want to spend most of "
          "them on what we could NOT do as much as what we could.")
p = pic(s, "01_cover.png", 0, 0, w=13.333)
if p:
    p.top = Inches(-0.9)
ov = s.shapes.add_shape(1, 0, Inches(4.5), SW, Inches(3.0))
ov.fill.solid(); ov.fill.fore_color.rgb = BG
ov.line.fill.background(); ov.shadow.inherit = False
tb(s, 0.7, 4.7, 12, 0.85, "Six Looks, One Harvest", size=46, bold=True)
tb(s, 0.7, 5.55, 12, 0.5,
   "A final yield forecast for 966 plots in Sokhda, Vadodara — from six Capella X-band passes",
   size=17, color=MUTE)
tb(s, 0.7, 6.3, 12, 0.4, "Team 8bit  ·  Harsh Thummar  ·  Viraj Suhagiya",
   size=16, bold=True, color=ACC)
tb(s, 0.7, 6.78, 12, 0.35,
   "ANRF AISEHack 2.0 · Round 3 · Remote Sensing: Yield Estimation", size=12.5, color=MUTE)

# ============================================================ 2 the problem
s = slide("Round 2 asked how the crop was doing on 13 October. Round 3 asks what comes off "
          "the field. Those are different problems. The six passes answer them very unevenly. "
          "Rice we see end to end. Cotton is still standing at the last pass, so its forecast "
          "is real extrapolation. Maize, bajra and groundnut have a sixty-day hole that "
          "swallows their entire maturity and harvest — one in-season look each. "
          "That asymmetry is not a caveat we bolted on. It sets the whole uncertainty structure.")
title(s, "Round 3 is a different problem",
      "Round 2 asked how the crop was doing. Round 3 asks what comes off the field.")
df = pd.DataFrame({
    "Crop": ["Rice", "Cotton", "Groundnut", "Maize", "Bajra"],
    "Cycle seen": ["100%", "78%", "60%", "55%", "50%"],
    "So the forecast is…": ["a reconstruction — season complete",
                            "genuine extrapolation past 12 Nov",
                            "one in-season look", "one in-season look",
                            "one in-season look"]})
table(s, df, 0.8, 1.6, 11.7, 2.7, fs=17, hdr=15, col_w=[2.2, 2.2, 7], hi=[0, 1])
tb(s, 0.8, 4.7, 11.7, 2.2,
   ["The 14 Aug → 13 Oct gap is **60 days**. It contains the entire maturity and harvest of "
    "maize, bajra and groundnut.",
    "",
    "**This is what sets the uncertainty structure — bajra's interval is widest, rice's is tightest.**"],
   size=18)

# ============================================================ 3 architecture
s = slide("This is the whole system on one slide. Left, what goes in. Centre, five stages "
          "from complex SLC to the forecast. Right, two things we tested and did not use — "
          "I will come back to those, because they matter. And bottom left, Sentinel-2, "
          "which is held out of every fitting step. That is the only reason our validation "
          "means anything.")
title(s, "System architecture", "Complex SLC in, at-harvest forecast out — and what we rejected on the way")
pic(s, "00_architecture.png", 0.25, 1.25, w=12.85)

# ============================================================ 4 measurement
s = slide("Everything from the complex SLC to the CSV is ours, built from the Capella "
          "metadata and the RPC model. No SNAP, no ISCE, no gdalwarp. Four numbers: our RPC "
          "reproduces the 225 control points Capella ships to four decimal places of a pixel. "
          "The ellipsoid-geoid datum is solved from those same control points, not assumed. "
          "Co-registration is within twelve centimetres on all six passes. And a plot mean is "
          "good to about a tenth of a decibel.")
title(s, "The measurement chain is ours, end to end",
      "Built from the Capella metadata and the RPC model — no SNAP, no ISCE, no gdalwarp")
stat(s, 0.5, 1.75, 3.0, "0.0000 px", "RPC geocoding error against\nthe 225 embedded GCPs")
stat(s, 3.6, 1.75, 3.0, "−62.02 m", "ellipsoid/geoid datum solved\nfrom the GCPs (sd 0.12 m)", GOOD)
stat(s, 6.7, 1.75, 3.0, "≤ 0.12 m", "inter-pass co-registration,\nall six passes")
stat(s, 9.8, 1.75, 3.0, "0.13 dB", "precision of the plot-mean γ⁰\non a 0.27 ha median plot")
tb(s, 0.8, 4.5, 11.7, 2.4,
   ["**β⁰ = (|DN| · scale_factor)²   →   γ⁰ = β⁰ · tan θ**",
    "",
    "γ⁰ because the six passes span 28.7°–35.2° incidence, and γ⁰ is the quantity that is, to "
    "first order, invariant to it. Local incidence is computed per range sample from the "
    "product state vectors."], size=17)

# ============================================================ 5 terrain
s = slide("Here is a mistake we made and caught. We first geocoded at a single constant "
          "height. All six passes agreed on minus twenty metres, residual under three metres. "
          "It looked fine. It was not — because a global shift metric recovers the mean "
          "alignment and completely hides terrain, and this block has twenty-six metres of "
          "relief, worth forty-six metres of ground range. The tell was the one right-looking "
          "pass. A height error pushes opposite look directions in opposite directions. Once "
          "we carried per-pixel terrain, that residual fell from 2.54 metres to 12 centimetres.")
title(s, "A correction we made to our own work",
      "26 m of relief becomes 46 m of ground-range error if you assume a flat Earth")
d2 = pd.DataFrame({
    "Check": ["RPC round-trip RMSE", "Correlation vs Capella's geocoding",
              "Co-registration, 29 Oct right-look", "Passes improved"],
    "Constant height": ["15.1–16.6 px", "0.769", "2.54 m", "—"],
    "Terrain-referenced": ["4.1–4.7 px", "0.840", "0.12 m", "6 of 6"]})
table(s, d2, 0.8, 1.6, 11.7, 2.3, fs=16, hdr=14, col_w=[5, 3, 3.4], hi=[0, 2])
tb(s, 0.8, 4.25, 11.7, 2.6,
   ["A height error displaces **opposite look directions in opposite senses**. So a residual "
    "that shows up on the single right-looking pass and nowhere else is a terrain signature, "
    "not a registration failure.",
    "",
    "**That one diagnostic is what exposed the error — and confirms the fix.**"], size=18)

# ============================================================ 6 geometry
s = slide("Geometry drove the design. Six passes, incidence from 28.7 to 35.2 degrees, one "
          "of them right-looking. No single normalisation removes that from a mixed scene. So "
          "two rules. First, every feature is a village anomaly — a plot minus the village "
          "median on that same date. In one step that cancels gain, incidence, look direction "
          "and the soil moisture common to all fields. Second, we only difference date pairs "
          "that are geometrically safe. There are exactly two.")
title(s, "Geometry drove the design, not the formula")
tb(s, 0.8, 1.55, 11.7, 1.5,
   ["**1 · Every feature is a village anomaly** — plot minus village median, same acquisition.",
    "Cancels calibration, gain, incidence response, look direction, and the common soil-moisture "
    "excursion, in one step and without fitting anything."], size=17.5)
d6 = pd.DataFrame({
    "Geometry-safe pair": ["19 Jun → 14 Aug", "13 Oct → 12 Nov"],
    "Δ incidence": ["0.08°", "1.78°"],
    "Look": ["same", "same"],
    "Measures": ["canopy establishment", "late retention vs harvest — new in Round 3"]})
table(s, d6, 0.8, 3.35, 11.7, 1.6, fs=16, hdr=14, col_w=[3, 2, 1.5, 5.5])
tb(s, 0.8, 5.35, 11.7, 1.6,
   ["**2 · Differences only between those two pairs.** The 29 Oct pass is right-looking from "
    "the opposite azimuth and is used only to corroborate, never on its own."], size=17.5)

# ============================================================ 7 x-band limits
s = slide("Now the uncomfortable measurement. On this village, the median change in gamma-"
          "nought at peak season is minus zero point five eight decibels relative to bare "
          "soil. Peak-season fields are darker than bare ground. That is physically right at "
          "HH — a closed canopy attenuates a rough tilled surface more than it scatters back. "
          "Combined with X-band's early saturation against LAI, it means inverting this data "
          "for biomass would be indefensible. So we do not. We use timing, retention and "
          "evenness instead.")
title(s, "What X-band HH can and cannot do here",
      "Measured on this village, not assumed from the literature")
tb(s, 0.8, 1.7, 11.7, 1.1,
   "**Median Δγ⁰ at peak season  =  −0.58 dB relative to bare soil**", size=30, color=BAD)
tb(s, 0.8, 3.0, 11.7, 3.9,
   ["Peak-season fields are *darker* than bare ground. At HH a developed canopy attenuates a "
    "rough, freshly-tilled soil surface more than its own volume return replaces it — and "
    "X-band saturates early against LAI.",
    "",
    "**So we do not invert X-band for biomass.** We use what a metre-scale X-band series "
    "genuinely resolves: phenological **timing**, late-season **retention**, and within-field "
    "**evenness**."], size=18)

# ============================================================ 8 the phase
s = slide("The question a SAR audience will ask: you had six SLCs, why only the amplitude? "
          "We answer it with numbers. Of fifteen pass pairs, five are opposite-look and nine "
          "are beyond the critical baseline — some by a factor of a hundred. Exactly one pair "
          "is geometrically viable, and it is 56 days apart, which is far beyond X-band "
          "coherence over a growing canopy. So repeat-pass InSAR is structurally unavailable. "
          "We then tried sub-look coherence, which is single-pass and immune to temporal "
          "decorrelation. The estimator works — point scatterers reach 0.98 against 0.25 for "
          "clutter. But at plot level it is null. We report it as a null.")
title(s, "You had six SLCs — why only the amplitude?",
      "Because we checked, rather than assumed")
d8 = pd.DataFrame({
    "Repeat-pass InSAR — all 15 pass pairs": ["Opposite look direction",
                                              "Beyond the critical baseline",
                                              "Geometrically viable"],
    "Pairs": ["5", "9", "1"],
    "": ["", "by up to 100×", "…and 56 days apart — no X-band coherence survives"]})
table(s, d8, 0.8, 1.55, 11.7, 1.9, fs=16, hdr=14, col_w=[5, 1.4, 5.6], hi=[2])
tb(s, 0.8, 3.75, 11.7, 3.2,
   ["**So we used the phase a different way.** Sub-aperture decomposition splits one "
    "acquisition's Doppler spectrum in two — single-pass, so temporal decorrelation cannot "
    "touch it.",
    "",
    "The estimator is sound: point scatterers reach **γ = 0.98** against **0.25** for distributed "
    "clutter. But at plot level it carries nothing — repeatability **r = 0.02**, crop separation "
    "**η² = 0.02**.",
    "",
    "**A null result, reported as a null.**"], size=17)

# ============================================================ 9 the model
s = slide("The model in one line. District statistics set the level; the SAR index sets the "
          "distribution around it. We separate those two deliberately, because with no ground "
          "truth the radar can rank plots but it cannot tell you the absolute level. The "
          "lognormal form makes the area-weighted village mean reproduce the anchor by "
          "construction — so the anchor is a stated assumption, not a hidden result.")
title(s, "The forecast model",
      "SAR sets the ranking and the spread. Official statistics set the level. Stated, not hidden.")
box = s.shapes.add_shape(1, Inches(0.8), Inches(1.55), Inches(11.7), Inches(1.15))
box.fill.solid(); box.fill.fore_color.rgb = BAND
box.line.fill.background(); box.shadow.inherit = False
tb(s, 0.95, 1.82, 11.4, 0.7,
   "Y_p  =  Y_ref(c) · S(c) · exp( σ_c · ρ · z_p  −  ½(σ_c · ρ)² )",
   size=27, bold=True, align=PP_ALIGN.CENTER, color=ACC)
d9 = pd.DataFrame({
    "Term": ["Y_ref(c)", "S(c)", "z_p", "σ_c , ρ"],
    "What it is": ["Vadodara district yield (DoA Gujarat 2022-23)",
                   "2025 season factor from the +16% rainfall anomaly",
                   "stage-weighted SAR index, standardised within crop",
                   "farm-to-farm yield CV, and the share of it SAR explains"],
    "Value": ["rice 1690 · cotton 776 · maize 2312 · bajra 2714 · groundnut 2514",
              "0.976 – 1.080, crop-specific",
              "establishment · Oct canopy · Nov retention · evenness",
              "0.28–0.35   ·   ρ = 0.55 (varied 0.40–0.70)"]})
table(s, d9, 0.8, 2.95, 11.7, 2.2, fs=13.5, hdr=13, col_w=[1.6, 5.0, 6.0])
tb(s, 0.8, 5.4, 11.7, 1.6,
   ["The lognormal form makes the **area-weighted village mean reproduce Y_ref · S by "
    "construction**. Cotton adds an explicit late-picking term — the one genuine forecast "
    "in the model."], size=16.5)

# ============================================================ 10 labels
s = slide("We re-derived the crop map from all six passes. It agrees with the carry-forward "
          "on only 46 per cent of plots. So we adjudicated using Sentinel-2, which neither "
          "classifier had seen. The carry-forward separates November greenness better on both "
          "dates. Our new classifier is not an improvement, so we did not adopt it. The "
          "disagreement becomes a measurement of label uncertainty instead. Crop labelling, "
          "not the radar, is the dominant uncertainty in this problem.")
title(s, "Crop labels: an honest negative result",
      "We re-derived the map from six passes. It was worse. So we did not use it.")
d10 = pd.DataFrame({
    "Label set": ["Round-1/2 carry-forward (4 passes)", "Our Round-3 re-derivation (6 passes)"],
    "η² · 12 Nov NDVI": ["0.138", "0.085"],
    "η² · 13 Oct NDVI": ["0.280", "0.217"],
    "Adopted": ["YES — primary", "no"]})
table(s, d10, 0.8, 1.6, 11.7, 1.7, fs=17, hdr=14, col_w=[5, 2.3, 2.3, 2.1], hi=[0])
tb(s, 0.8, 3.7, 11.7, 3.2,
   ["The two sets agree on only **46%** of observed plots. Adjudicated on data neither had "
    "seen, the carry-forward wins on both dates.",
    "",
    "**The disagreement is not discarded — it becomes a measured label uncertainty**, carried "
    "through to the village result.",
    "",
    "**Crop labelling, not the radar, is the dominant uncertainty here.**"], size=18)

# ============================================================ 11 validation
s = slide("Validation. Two Sentinel-2 scenes fall on exactly the same day as Capella passes, "
          "so no interpolation. The headline is the within-crop result. Our yield index "
          "correlates significantly for rice and cotton — the two crops still standing on 13 "
          "October — and is statistically null for maize, bajra and groundnut, the three "
          "already cut, where both sensors see bare soil. A speckle artefact would correlate "
          "everywhere or nowhere. This validates exactly where a crop is still in the field. "
          "And on the right, why SAR: a 111-day optical blackout across the whole monsoon.")
title(s, "The index works exactly where the crop is still standing",
      "Sentinel-2 tests the result — it never enters the model")
d11 = pd.DataFrame({
    "Crop": ["Rice", "Cotton", "Groundnut", "Maize", "Bajra"],
    "ρ vs same-day NDVI": ["+0.365", "+0.361", "+0.112", "+0.019", "−0.071"],
    "p": ["2×10⁻⁸", "1×10⁻¹²", "0.079", "0.88", "0.61"],
    "On 13 Oct": ["standing", "standing", "harvesting", "harvested", "harvested"]})
table(s, d11, 0.6, 1.55, 7.3, 2.4, fs=14.5, hdr=12.5, col_w=[1.8, 3.0, 1.8, 2.4], hi=[0, 1])
tb(s, 0.6, 4.25, 7.3, 2.7,
   ["**Significant for the two crops still in the field. Null for the three already cut.**",
    "",
    "An artefact would correlate everywhere or nowhere. This pattern is hard to produce by "
    "accident.",
    "",
    "Whole village, same day: ρ = +0.56 (13 Oct) · +0.48 (12 Nov)"], size=15)
pic(s, "07_optical_blackout.png", 8.2, 1.9, w=4.7)
tb(s, 8.2, 5.15, 4.7, 1.8,
   ["**Why it must be SAR.** 46 Sentinel-2 overpasses Jun–Nov, 10 usable. Between 19 Jun and "
    "8 Oct: 29 passes, none below 23.3% cloud. A **111-day blackout** over sowing, "
    "establishment and peak growth. All six SAR passes clean."], size=13)

# ============================================================ 12 map
s = slide("The plot-level forecast, all 966 plots. What I want you to notice is that the "
          "pattern forms coherent spatial blocks — contiguous stronger and weaker areas — "
          "rather than plot-level noise. That is what genuine agronomic variation looks like. "
          "A speckle-driven artefact would not do that.")
title(s, "Plot-level final yield forecast — all 966 plots",
      "Coherent spatial structure, not plot-level speckle")
pic(s, "02_yield_and_crop_map.png", 0.3, 1.3, w=12.75)

# ============================================================ 13 results
s = slide("The village result. 756 tonnes across 447 hectares. Every village mean sits on its "
          "district anchor by construction; the intervals come from the Monte Carlo. Note that "
          "the interval widths order correctly with how much of each crop we actually saw — "
          "bajra widest, rice tightest. Cotton is reported as lint; that is about 2,340 "
          "kilograms per hectare of seed cotton.")
title(s, "Village-level forecast — Sokhda, 447.5 ha")
v = vil.set_index("crop_type").loc[["Rice", "Cotton", "Maize", "Bajra", "Groundnut"]]
d13 = pd.DataFrame({
    "Crop": v.index,
    "Plots": v.n_plots.values,
    "Area (ha)": [f"{x:.1f}" for x in v.area_ha],
    "Forecast (kg/ha)": [f"{x:,.0f}" for x in v.yield_forecast_kg_ha],
    "P10 – P90": [f"{a:,.0f} – {b:,.0f}" for a, b in
                  zip(v.yield_p10_kg_ha, v.yield_p90_kg_ha)],
    "Production (t)": [f"{x:.1f}" for x in v.production_t],
    "Product": v["product"].values})
table(s, d13, 0.5, 1.5, 12.35, 2.7, fs=15, hdr=13,
      col_w=[2.0, 1.2, 1.6, 2.2, 2.5, 2.0, 2.2])
stat(s, 0.5, 4.6, 3.9, f"{vil.production_t.sum():.0f} t", "total forecast production", bs=38)
stat(s, 4.7, 4.6, 3.9, "966 / 966", "plots carry a forecast\n832 full · 91 partial · 43 flagged", bs=38)
stat(s, 8.9, 4.6, 3.9, "±18–20%", "village interval, ordered by\nhow much of the crop we saw", WARN, bs=38)
tb(s, 0.5, 6.75, 12.35, 0.5,
   "Area-weighted, never a plain plot mean — plot areas span three orders of magnitude. "
   "Cotton as lint; 820 kg/ha ≈ 2,340 kg/ha seed cotton.", size=13, color=MUTE)

# ============================================================ 14 crop mix
s = slide("And here is the honest headline risk. The carry-forward puts groundnut at 28 per "
          "cent of Sokhda. The Directorate of Agriculture puts groundnut at nought point "
          "three five per cent of the district. Those cannot both be right. We ran the "
          "assignment under both area constraints and publish both. Village production moves "
          "from 763 tonnes to 570 — a 25 per cent swing, essentially all of it groundnut. "
          "Our plot-level radar work is solid. The crop map we inherited is the thing that "
          "would move the answer, and we would rather quantify that than hide it.")
title(s, "The one number that would move the answer",
      "Not the radar — the crop map we inherited")
d14 = pd.DataFrame({
    "Crop": ["Rice", "Cotton", "Maize", "Bajra", "Groundnut", "VILLAGE TOTAL"],
    "A · Round-1 carry-forward": ["61 ha → 106 t", "188 ha → 154 t", "35 ha → 88 t",
                                  "35 ha → 103 t", "127 ha → 313 t", "763 t"],
    "B · District statistics": ["79 ha → 135 t", "292 ha → 239 t", "64 ha → 159 t",
                               "11 ha → 32 t", "1.7 ha → 4 t", "570 t"]})
table(s, d14, 0.8, 1.55, 11.7, 2.9, fs=15, hdr=14, col_w=[2.4, 4.6, 4.6], hi=[4, 5])
tb(s, 0.8, 4.75, 11.7, 2.2,
   ["Round-1 puts groundnut at **28% of Sokhda**; the district figure is **0.35%**. "
    "Both cannot be close to right.",
    "",
    "**A 25% swing in village production, essentially all of it groundnut.** We publish both "
    "rather than choose silently."], size=18)

# ============================================================ 15 limits
s = slide("Where this is weakest, plainly. No ground truth exists, so absolute accuracy is "
          "untestable — only internal consistency and independent optical agreement. Three "
          "crops rest on a single in-season look. Labels dominate the error budget. Only five "
          "classes are allowed but Vadodara also grows castor and pigeon pea at areas like "
          "maize, so some plots are certainly something else. And 43 plots fall outside every "
          "swath — flagged, not hidden.")
title(s, "Where this is weakest",
      "A forecast without its limits is not a forecast")
tb(s, 0.8, 1.6, 11.7, 5.4,
   ["**No ground truth exists** — absolute accuracy is untestable. Only internal consistency "
    "and independent optical agreement can be checked, and both are reported.",
    "",
    "**Maize, bajra and groundnut rest on one in-season observation.** Their intervals are "
    "wide by construction and should be read that way.",
    "",
    "**Crop labels dominate the error budget, not the radar** — a 25% swing in village "
    "production.",
    "",
    "**Only five classes are permitted**, but Vadodara also grows castor and pigeon pea at "
    "areas comparable to maize. Some plots are certainly another crop.",
    "",
    "**43 plots fall outside every swath footprint** — they carry the crop mean with inflated "
    "uncertainty and are flagged, not hidden."], size=17, space_after=0)

# ============================================================ 16 close
s = slide("What we would ask for next. One acquisition in early September would convert three "
          "crops from one-look estimates into properly sampled seasons — the biggest gain "
          "available for the least data. Any dual-pol acquisition would let us separate volume "
          "from surface scattering, which is what forced us away from biomass inversion. And "
          "twenty ground-truth plots — not to fit the model, but to estimate the one parameter "
          "we currently have to assume. Thank you.")
bar = s.shapes.add_shape(1, 0, 0, SW, SH)
bar.fill.solid(); bar.fill.fore_color.rgb = BAND
bar.line.fill.background(); bar.shadow.inherit = False
tb(s, 1.0, 1.3, 11.3, 0.9, "What we would want next", size=36, bold=True)
tb(s, 1.0, 2.5, 11.3, 3.8,
   ["**One acquisition in early September.** A single pass inside the 60-day gap would convert "
    "maize, bajra and groundnut from one-look estimates into properly sampled seasons — the "
    "largest gain available for the least data.",
    "",
    "**Any dual-polarisation acquisition.** HH alone cannot separate volume from surface "
    "scattering; HV would resolve much of what forced us away from biomass inversion.",
    "",
    "**Twenty ground-truth plots.** Not to fit the model — to estimate ρ, the one parameter we "
    "currently have to assume."], size=18)
tb(s, 1.0, 6.6, 11.3, 0.5,
   "Team 8bit  ·  Harsh Thummar, Viraj Suhagiya  ·  ANRF AISEHack 2.0 Round 3",
   size=14, color=MUTE)

path = os.path.join(DEL, "Team8bit_Round3_GoaFinals.pptx")
prs.save(path)
n = len(prs.slides._sldIdLst)
print(f"wrote {path}  ({os.path.getsize(path)/1024:.0f} KB, {n} slides)")
print("speaker notes attached to every slide")
